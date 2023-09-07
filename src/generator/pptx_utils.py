"""
Provides a set of functions that directly interact with the PowerPoint presentation being generated.

Authors: Will Hoover, Cade Reinberger
"""
import copy
from pptx.util import Inches

def duplicate_question_slide(slide, pres):
    """
    Adds and returns a duplicate of the last slide in the given presentation.
    Written by Cade Reinberger
    """
    template = slide
    blank_slide_layout = pres.slide_layouts[5]

    copied_slide = pres.slides.add_slide(blank_slide_layout)
    copied_slide.title.text = template.title.text
    for i in range (1, len(template.shapes)):
        el = template.shapes[i].element
        newel = copy.deepcopy(copy.deepcopy(el))
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
    copied_slide.name = f'{slide}'

    return copied_slide

def add_new_question(pres, header: str):
    """
    Adds a new slide that is the start of the next tossup or bonus. This slide will be duplicated to
    create further slides for this question.
    Returns the slide that was created.
    """
    new_slide = pres.slides.add_slide(pres.slide_master.slide_layouts[5])
    new_slide.title.text = header
    new_slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(12), Inches(7))
    return new_slide

def add_question_fragment(prev_slide, pres, fragment: str):
    """
    Creates and adds a new slide with the given question fragment appended to it.
    Returns the slide that was created.
    """
    new_slide = duplicate_question_slide(prev_slide, pres)
    # fragment_run = new_slide.shapes[1].text_frame.paragraphs[0].add_run()
    # fragment_run.text = fragment
    return new_slide

def add_tossup_answer(prev_slide, pres, answer: str):
    """
    Creates and adds a new slide with the given answer to the finished tossup inserted at the bottom.
    Returns the slide that was created.
    """
    return add_question_fragment(prev_slide, pres, answer)